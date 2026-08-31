import os
import re
import json
import logging
from dataclasses import dataclass
from typing import List, Tuple
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

logger = logging.getLogger(__name__)


@dataclass
class PolicyChunk:
    doc_name: str
    section_title: str
    content: str

    @property
    def source_reference(self) -> str:
        return f"{self.doc_name} > {self.section_title}"


class PolicyRetriever:
    """
    Lightweight policy document loader, multi-format parser (.pdf, .pptx, .json, .md, .txt),
    text chunker, and TF-IDF cosine similarity retriever.
    """

    ALLOWED_EXTENSIONS = {".md", ".txt", ".pdf", ".pptx", ".ppt", ".json"}

    def __init__(self, policy_dir: str):
        self.policy_dir = policy_dir
        self.chunks: List[PolicyChunk] = []
        self.vectorizer: TfidfVectorizer = TfidfVectorizer(stop_words="english", ngram_range=(1, 2), sublinear_tf=True)
        self.tfidf_matrix = None
        self._load_and_chunk_docs()

    def _load_and_chunk_docs(self):
        """Scans policy_dir, reads all supported document formats, and chunks them."""
        if not os.path.exists(self.policy_dir):
            os.makedirs(self.policy_dir, exist_ok=True)

        for filename in sorted(os.listdir(self.policy_dir)):
            ext = os.path.splitext(filename)[1].lower()
            if ext in self.ALLOWED_EXTENSIONS:
                filepath = os.path.join(self.policy_dir, filename)
                try:
                    self._parse_and_chunk_file(filename, filepath, ext)
                except Exception as e:
                    logger.warning(f"Error parsing document '{filename}': {e}")

        if not self.chunks:
            # Fallback default empty chunk if directory is empty
            self.chunks.append(PolicyChunk(
                doc_name="General Policy",
                section_title="Overview",
                content="IIFL Finance Policy Knowledge Base."
            ))

        # Build TF-IDF matrix over chunk contents with document metadata
        corpus = [f"{c.doc_name} {c.section_title} {c.doc_name}\n{c.content}" for c in self.chunks]
        self.vectorizer = TfidfVectorizer(stop_words="english", ngram_range=(1, 2), sublinear_tf=True)
        self.tfidf_matrix = self.vectorizer.fit_transform(corpus)

    def _parse_and_chunk_file(self, filename: str, filepath: str, ext: str):
        """Dispatches file parsing according to its extension."""
        doc_title = os.path.splitext(filename)[0].replace("_", " ").title()

        if ext in [".md", ".txt"]:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            self._chunk_markdown(doc_title, text)

        elif ext == ".pdf":
            self._parse_pdf(doc_title, filepath)

        elif ext in [".pptx", ".ppt"]:
            self._parse_pptx(doc_title, filepath)

        elif ext == ".json":
            self._parse_json(doc_title, filepath)

    def _chunk_markdown(self, doc_title: str, text: str):
        """Splits markdown text by section headers (#, ##, ###)."""
        sections = re.split(r'\n(?=#{1,3}\s+)', text)
        for section in sections:
            section = section.strip()
            if not section:
                continue
            lines = section.split('\n')
            first_line = lines[0].strip()
            if first_line.startswith('#'):
                section_title = re.sub(r'^#{1,3}\s+', '', first_line).strip()
                content = '\n'.join(lines[1:]).strip() if len(lines) > 1 else section
            else:
                section_title = "General Information"
                content = section

            if content:
                self.chunks.append(PolicyChunk(
                    doc_name=doc_title,
                    section_title=section_title,
                    content=content
                ))

    def _parse_pdf(self, doc_title: str, filepath: str):
        """Extracts text from PDF documents page by page."""
        try:
            from pypdf import PdfReader
            reader = PdfReader(filepath)
            for page_idx, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                page_text = page_text.strip()
                if page_text:
                    self.chunks.append(PolicyChunk(
                        doc_name=doc_title,
                        section_title=f"Page {page_idx + 1}",
                        content=page_text
                    ))
        except Exception as e:
            logger.warning(f"pypdf extraction failed for {filepath}: {e}")

    def _parse_pptx(self, doc_title: str, filepath: str):
        """Extracts text from PowerPoint presentations slide by slide."""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            txt = paragraph.text.strip()
                            if txt:
                                slide_texts.append(txt)
                
                content = "\n".join(slide_texts).strip()
                if content:
                    # Try to use first line as slide title
                    slide_title = slide_texts[0] if slide_texts else f"Slide {slide_idx + 1}"
                    if len(slide_title) > 60:
                        slide_title = f"Slide {slide_idx + 1}"
                    self.chunks.append(PolicyChunk(
                        doc_name=doc_title,
                        section_title=f"Slide {slide_idx + 1} ({slide_title})",
                        content=content
                    ))
        except Exception as e:
            logger.warning(f"pptx extraction failed for {filepath}: {e}")

    def _parse_json(self, doc_title: str, filepath: str):
        """Extracts structured text from JSON files."""
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)

            if isinstance(data, list):
                for idx, item in enumerate(data):
                    title = item.get("title") or item.get("topic") or item.get("category") or f"Item {idx + 1}"
                    content = json.dumps(item, indent=2)
                    self.chunks.append(PolicyChunk(
                        doc_name=doc_title,
                        section_title=str(title),
                        content=content
                    ))
            elif isinstance(data, dict):
                for key, val in data.items():
                    content = json.dumps(val, indent=2) if isinstance(val, (dict, list)) else str(val)
                    self.chunks.append(PolicyChunk(
                        doc_name=doc_title,
                        section_title=str(key),
                        content=content
                    ))
        except Exception as e:
            logger.warning(f"JSON extraction failed for {filepath}: {e}")

    def retrieve(self, query: str, top_k: int = 5, score_threshold: float = 0.04) -> List[Tuple[PolicyChunk, float]]:
        """
        Retrieves top_k matching chunks for a given query string.
        Returns a list of (PolicyChunk, similarity_score) tuples.
        """
        if not query or not query.strip() or self.tfidf_matrix is None:
            return []

        query_vec = self.vectorizer.transform([query])
        similarities = cosine_similarity(query_vec, self.tfidf_matrix).flatten()

        top_indices = similarities.argsort()[::-1][:top_k]
        results = []

        for idx in top_indices:
            score = float(similarities[idx])
            if score >= score_threshold:
                results.append((self.chunks[idx], score))

        return results
